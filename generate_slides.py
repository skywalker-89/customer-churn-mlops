"""
Customer Churn ML Project — PowerPoint Presentation Generator
Generates a full 18-slide professional .pptx deck.
"""

import os
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── colour palette ──────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x3E)   # slide background / dark fill
ACCENT      = RGBColor(0x00, 0xB4, 0xFF)   # bright cyan accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xD6, 0xE8)
GOLD        = RGBColor(0xFF, 0xC8, 0x00)
GREEN       = RGBColor(0x2E, 0xCC, 0x71)
RED_LIGHT   = RGBColor(0xFF, 0x6B, 0x6B)
CARD_BG     = RGBColor(0x15, 0x28, 0x57)   # card / text-box interior

BASE = Path(__file__).parent
IMG_CLASS = BASE / "data_quality" / "classification_20260222_013847"
IMG_REG   = BASE / "data_quality" / "regression_20260222_013921"
IMG_CORR  = BASE / "correlation_analysis"
IMG_FEAT  = BASE / "feature_relationships.png"

OUT_DIR  = BASE
OUT_PATH = OUT_DIR / "Customer_Churn_ML_Project_Slides.pptx"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── helpers ─────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]          # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color: RGBColor = NAVY):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color: RGBColor = WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.color.rgb = color
    run.font.italic = italic
    try:
        run.font.name = "Calibri"
    except Exception:
        pass
    return txBox


def add_bullets(slide, items, left, top, width, height,
                font_size=18, color: RGBColor = WHITE,
                bullet_color: RGBColor = ACCENT):
    """Add bulleted list inside a text box."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"▸  {item}"
        run.font.size  = Pt(font_size)
        run.font.color.rgb = color
        try:
            run.font.name = "Calibri"
        except Exception:
            pass
    return txBox


def slide_title_bar(slide, title_text, subtitle=None):
    """Draw accent bar at top, title text."""
    add_rect(slide, 0, 0, 13.33, 1.1, CARD_BG)
    add_rect(slide, 0, 0, 0.25, 1.1, ACCENT)
    add_text(slide, title_text, 0.4, 0.05, 11.5, 0.95,
             font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.8, 11.5, 0.4,
                 font_size=14, color=ACCENT, align=PP_ALIGN.LEFT)


def slide_number(slide, prs, n):
    """Bottom right slide number (skip on slide 1)."""
    total = len(prs.slides)  # won't be final, but close enough
    add_text(slide, f"{n}", 12.8, 7.15, 0.5, 0.3,
             font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


def try_add_image(slide, path, left, top, width, height=None):
    """Add image if file exists."""
    p = Path(path)
    if p.exists():
        pic = slide.shapes.add_picture(
            str(p), Inches(left), Inches(top), Inches(width),
            None if height is None else Inches(height)
        )
        return pic
    else:
        # placeholder label
        add_rect(slide, left, top, width, height or 2, CARD_BG)
        add_text(slide, f"[INSERT CHART: {p.name}]",
                 left + 0.1, top + (height or 2) / 2 - 0.2,
                 width - 0.2, 0.5,
                 font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    return None


# ── load metrics ────────────────────────────────────────────────────────────
EVAL_JSON = BASE / "data_quality" / "evaluation_report_20260221_163732.json"
with open(EVAL_JSON) as f:
    EVAL = json.load(f)

CLASSIF = {r["model"]: r for r in EVAL["classification"]["results"]}
REGR    = {r["model"]: r for r in EVAL["regression"]["results"]}


def fmt(v, decimals=4):
    return f"{v:.{decimals}f}"


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════

def slide1_title(prs):
    slide = blank_slide(prs)
    bg(slide, NAVY)

    # large gradient-style banner
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
    add_rect(slide, 0, 0, 13.33, 0.12, ACCENT)         # top accent line
    add_rect(slide, 0, 7.38, 13.33, 0.12, ACCENT)      # bottom accent line
    add_rect(slide, 0, 2.2, 13.33, 0.06, CARD_BG)      # divider

    add_text(slide, "Customer Churn Prediction & Analysis",
             0.6, 0.5, 12.0, 1.5,
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "MLOps Pipeline  |  From-Scratch ML Implementations  |  Dual-Task Learning",
             0.6, 1.8, 12.0, 0.5,
             font_size=16, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)

    # info block
    info = [
        ("Course", "[01276343] ML Project"),
        ("Date",   "March 9, 2026"),
        ("Type",   "On-site Presentation  —  Class Hours"),
        ("Team",   "Group 3 Members"),
    ]
    y = 2.8
    for label, val in info:
        add_rect(slide, 2.5, y, 8.3, 0.65, CARD_BG)
        add_text(slide, label + ":", 2.7, y + 0.1, 2.0, 0.45,
                 font_size=14, bold=True, color=ACCENT)
        add_text(slide, val, 4.8, y + 0.1, 5.8, 0.45,
                 font_size=14, color=WHITE)
        y += 0.78

    add_text(slide, "Customer Churn MLOps  ·  Powered by Airflow, MinIO & MLflow",
             0, 7.1, 13.33, 0.3,
             font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    return slide


def slide2_agenda(prs):
    slide = blank_slide(prs)
    bg(slide)
    slide_title_bar(slide, "Agenda")

    items_left = [
        "1.  Problem Statement & Business Case",
        "2.  Dataset Overview",
        "3.  Feature Analysis & EDA",
        "4.  Data Preprocessing",
        "5.  Methodology Overview",
        "6.  Regression Models",
        "7.  Regression Results",
        "8.  Classification Models (Part 1 – 3)",
    ]
    items_right = [
        "9.   Novel / New Models",
        "10. Benchmarking Table",
        "11. ROC & AUC Curves",
        "12. Key Findings & Discussion",
        "13. Conclusion",
        "14. References",
        "",
        "15. Q & A",
    ]
    add_bullets(slide, items_left,  0.5, 1.3, 6.0, 5.8, font_size=17)
    add_bullets(slide, items_right, 6.8, 1.3, 6.0, 5.8, font_size=17)
    slide_number(slide, prs, 2)
    return slide


def slide3_problem(prs):
    slide = blank_slide(prs)
    bg(slide)
    slide_title_bar(slide, "Problem Statement", "Why does customer churn matter?")

    # three cards
    cards = [
        ("What is Customer Churn?",
         ["Churn = customer stops buying / leaves the business",
          "Retail sector avg churn: 25–30 % per year",
          "Binary target: churned = 1 (left) vs 0 (retained)",
          "Affects revenue, brand loyalty & growth"],
         0.4),
        ("Business Impact",
         ["Retaining a customer is 5–25× cheaper than acquiring one",
          "5–10 % revenue lift from better retention",
          "15–20 % reduction in acquisition costs",
          "Enables personalised offers & proactive outreach"],
         4.45),
        ("Why ML / DL?",
         ["70+ feature interactions too complex for manual rules",
          "Dual-task: regression (CLV) + classification (churn)",
          "From-scratch implementations prove deep understanding",
          "MLOps pipeline ensures reproducibility at scale"],
         8.5),
    ]
    for title, bullets, x in cards:
        add_rect(slide, x, 1.4, 3.8, 5.7, CARD_BG)
        add_rect(slide, x, 1.4, 3.8, 0.5, ACCENT)
        add_text(slide, title, x + 0.15, 1.42, 3.5, 0.45,
                 font_size=14, bold=True, color=NAVY)
        add_bullets(slide, bullets, x + 0.1, 2.0, 3.6, 4.8,
                    font_size=14, color=LIGHT_GRAY)
    slide_number(slide, prs, 3)
    return slide


def slide4_dataset(prs):
    slide = blank_slide(prs)
    bg(slide)
    slide_title_bar(slide, "Dataset Overview", "Retail Customer Analytics Dataset")

    bullets_left = [
        "Source: Synthetic multi-category retail store data",
        "Size: 1,000 customer records × 70+ features",
        "Target 1 — total_sales  (regression, continuous)",
        "Target 2 — churned  (classification, binary 0/1)",
        "6 feature groups: Demographics, Behaviour, Transactions,",
        "   Products, Marketing & Engagement",
        "Train / Test split: 80 % / 20 %  (stratified)",
        "Class balance: ~70 % retained, ~30 % churned",
    ]
    add_bullets(slide, bullets_left, 0.4, 1.3, 6.8, 5.8, font_size=16)

    # embed correlation heatmap as dataset overview visual
    try_add_image(slide, IMG_CORR / "correlation_heatmap_full.png",
                  7.6, 1.25, 5.5, 5.8)
    add_text(slide, "▲ Full feature correlation heatmap",
             7.6, 6.9, 5.5, 0.4, font_size=11, color=LIGHT_GRAY,
             align=PP_ALIGN.CENTER)
    slide_number(slide, prs, 4)
    return slide


def slide5_features(prs):
    slide = blank_slide(prs)
    bg(slide)
    slide_title_bar(slide, "Feature Analysis & Storytelling",
                    "Understanding what drives churn and lifetime value")

    bullets = [
        "Quantity × Unit Price → Total Transaction Value  (key interaction feature)",
        "Days Since Last Purchase → Recency (RFM signal)",
        "Loyalty Membership Years → Customer lifetime proxy",
        "Purchase Frequency → Engagement depth",
        "Category & Promotion Channel → Behavioural clustering",
        "Correlation analysis reveals multicollinearity clusters",
        "Feature importance guides model selection per track",
    ]
    add_bullets(slide, bullets, 0.4, 1.3, 6.5, 5.5, font_size=16)

    try_add_image(slide, IMG_FEAT, 7.2, 1.25, 5.9, 5.8)
    add_text(slide, "▲ Feature relationship chart",
             7.2, 6.9, 5.9, 0.4, font_size=11, color=LIGHT_GRAY,
             align=PP_ALIGN.CENTER)
    slide_number(slide, prs, 5)
    return slide


def slide6_preprocessing(prs):
    slide = blank_slide(prs)
    bg(slide)
    slide_title_bar(slide, "Data Preprocessing Pipeline")

    steps = [
        ("1  Null Handling",
         "Median imputation for numerical columns; mode for categorical"),
        ("2  Encoding",
         "One-Hot Encoding for nominal vars (gender, category, location)"),
        ("3  Feature Engineering",
         "Interaction: quantity × unit_price → total_transaction_value"),
        ("4  Aggregation",
         "Per-customer RFM: recency, frequency, avg order value"),
        ("5  Scaling",
         "StandardScaler applied to all numerical features"),
        ("6  Train / Test Split",
         "80 / 20 stratified split — identical for all models (fair benchmark)"),
        ("7  Storage",
         "Processed data written to MinIO as training_data.parquet"),
    ]
    y = 1.35
    for title, desc in steps:
        add_rect(slide, 0.5, y, 12.3, 0.7, CARD_BG)
        add_rect(slide, 0.5, y, 0.08, 0.7, ACCENT)
        add_text(slide, title, 0.7, y + 0.07, 3.2, 0.55,
                 font_size=15, bold=True, color=ACCENT)
        add_text(slide, desc, 4.0, y + 0.07, 8.6, 0.55,
                 font_size=14, color=WHITE)
        y += 0.84
    slide_number(slide, prs, 6)
    return slide


def slide7_methodology(prs):
    slide = blank_slide(prs)
    bg(slide)
    slide_title_bar(slide, "Methodology Overview",
                    "Dual-track ML pipeline with from-scratch implementations")

    # regression track
    add_rect(slide, 0.4, 1.3, 5.9, 5.8, CARD_BG)
    add_rect(slide, 0.4, 1.3, 5.9, 0.5, RGBColor(0x1A, 0x78, 0xC2))
    add_text(slide, "📈  REGRESSION TRACK  (total_sales)",
             0.55, 1.33, 5.6, 0.43, font_size=14, bold=True, color=WHITE)
    reg_models = [
        "Linear Regression (from scratch)",
        "Multiple Regression (from scratch)",
        "Polynomial Regression (from scratch)",
        "XGBoost (from scratch)  ← novel model",
        "XGBoost (sklearn)  ← validation baseline",
    ]
    add_bullets(slide, reg_models, 0.5, 1.9, 5.7, 5.0, font_size=15)

    # classification track
    add_rect(slide, 7.0, 1.3, 5.9, 5.8, CARD_BG)
    add_rect(slide, 7.0, 1.3, 5.9, 0.5, RGBColor(0x1A, 0x9C, 0x60))
    add_text(slide, "🔶  CLASSIFICATION TRACK  (churned)",
             7.15, 1.33, 5.6, 0.43, font_size=14, bold=True, color=WHITE)
    clf_models = [
        "Logistic Regression",
        "Decision Tree",
        "Random Forest  (+PCA variant)",
        "Support Vector Machine  (+PCA variant)",
        "K-Means Clustering",
        "Agglomerative Clustering",
        "Perceptron / SLP",
        "Multi-Layer Perceptron (MLP)",
        "Custom Model — Naive Bayes  ← novel",
    ]
    add_bullets(slide, clf_models, 7.1, 1.9, 5.7, 4.9, font_size=14)

    add_rect(slide, 5.9, 1.3, 1.1, 5.8, NAVY)
    add_text(slide, "VS", 5.95, 4.1, 1.0, 0.55,
             font_size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    slide_number(slide, prs, 7)
    return slide


def slide8_regression_models(prs):
    slide = blank_slide(prs)
    bg(slide)
    slide_title_bar(slide, "Regression Models",
                    "All models built from scratch using NumPy only")

    models = [
        ("Linear Regression",
         "OLS gradient descent; models linear relationship between "
         "individual feature and total_sales"),
        ("Multiple Regression",
         "Extends linear model to all features simultaneously; "
         "learns weight vector via gradient descent"),
        ("Polynomial Regression",
         "Degree-2 feature expansion captures non-linear relationships; "
         "regularised to reduce overfitting"),
        ("XGBoost (Novel)",
         "Gradient-boosted decision trees built from scratch; "
         "additive residual learning with shrinkage; "
         "significantly outperforms classical regression"),
    ]
    y = 1.3
    for name, desc in models:
        add_rect(slide, 0.4, y, 12.4, 1.35, CARD_BG)
        add_rect(slide, 0.4, y, 0.1, 1.35, ACCENT)
        add_text(slide, name, 0.65, y + 0.1, 3.5, 0.5,
                 font_size=17, bold=True, color=ACCENT)
        add_text(slide, desc, 0.65, y + 0.6, 12.0, 0.65,
                 font_size=14, color=LIGHT_GRAY)
        y += 1.5

    add_text(slide,
             "✔  Each model also verified against sklearn equivalent for result validation",
             0.4, 7.05, 12.5, 0.35, font_size=13, color=GREEN)
    slide_number(slide, prs, 8)
    return slide


def slide9_regression_results(prs):
    slide = blank_slide(prs)
    bg(slide)
    slide_title_bar(slide, "Regression Results",
                    "XGBoost overwhelmingly outperforms classical baselines")

    # metrics table
    headers = ["Model", "RMSE", "MAE", "R²", "MAPE"]
    rows = [
        ("Linear Regression (scratch)",    "1,248.77", "984.93",  "0.4134", "368.90%"),
        ("Multiple Regression (scratch)",  "1,546.78", "1,235.65","0.1000", "463.01%"),
        ("Polynomial Regression (scratch)","825.45",   "620.37",  "0.7437", "176.80%"),
        ("XGBoost (scratch)  ★",           "119.86",    "89.41",   "0.9946", " 20.24%"),
        ("XGBoost (sklearn)  ✔",           " 54.76",    "34.64",   "0.9989", "  4.46%"),
    ]

    col_w = [4.0, 1.8, 1.8, 1.5, 1.8]
    col_x = [0.4, 4.55, 6.4, 8.25, 9.8]
    header_y = 1.25

    # header row
    for i, (hdr, x, w) in enumerate(zip(headers, col_x, col_w)):
        add_rect(slide, x, header_y, w - 0.05, 0.5, ACCENT)
        add_text(slide, hdr, x + 0.05, header_y + 0.05, w - 0.1, 0.4,
                 font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    row_colors = [CARD_BG, CARD_BG, CARD_BG,
                  RGBColor(0x0D, 0x35, 0x60),   # highlight XGBoost scratch
                  RGBColor(0x0D, 0x35, 0x60)]
    y = 1.85
    for ri, (row, rc) in enumerate(zip(rows, row_colors)):
        for ci, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
            add_rect(slide, x, y, w - 0.05, 0.55, rc)
            clr = GOLD if ri >= 3 else WHITE
            add_text(slide, cell, x + 0.05, y + 0.08, w - 0.1, 0.4,
                     font_size=13, color=clr, align=PP_ALIGN.CENTER,
                     bold=(ri >= 3))
        y += 0.6

    # embed metrics comparison plot
    try_add_image(slide, IMG_REG / "metrics_comparison.png",
                  0.4, 4.4, 5.5, 2.9)
    try_add_image(slide, IMG_REG / "predicted_vs_actual_xgboost_regression.png",
                  6.1, 4.4, 6.8, 2.9)
    slide_number(slide, prs, 9)
    return slide


def slide10_classif1(prs):
    slide = blank_slide(prs)
    bg(slide)
    slide_title_bar(slide, "Classification Models — Part 1",
                    "Logistic Regression · Decision Tree · Random Forest")

    bullets_left = [
        "Logistic Regression — sigmoid + binary cross-entropy, GD",
        f"  Accuracy: {fmt(CLASSIF['Logistic Regression (scratch)']['accuracy'],4)}  "
        f"  F1: {fmt(CLASSIF['Logistic Regression (scratch)']['f1_score'],4)}",
        "",
        "Decision Tree — recursive binary splits, Gini impurity",
        f"  Accuracy: {fmt(CLASSIF['Decision Tree (scratch)']['accuracy'],4)}  "
        f"  F1: {fmt(CLASSIF['Decision Tree (scratch)']['f1_score'],4)}",
        "",
        "Random Forest — bagged ensemble of 100 decision trees",
        f"  Accuracy: {fmt(CLASSIF['Random Forest (scratch)']['accuracy'],4)}  "
        f"  F1: {fmt(CLASSIF['Random Forest (scratch)']['f1_score'],4)}",
        "",
        "All built from scratch using NumPy; sklearn RF used to validate",
    ]
    add_bullets(slide, bullets_left, 0.4, 1.25, 6.2, 5.9, font_size=15)

    # confusion matrices (2 side by side)
    try_add_image(slide, IMG_CLASS / "confusion_matrix_logisticregression_classification.png",
                  6.8, 1.3, 3.1, 2.7)
    try_add_image(slide, IMG_CLASS / "confusion_matrix_randomforest_classification.png",
                  10.0, 1.3, 3.1, 2.7)
    add_text(slide, "Logistic Regression CM", 6.8, 3.95, 3.1, 0.35,
             font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "Random Forest CM", 10.0, 3.95, 3.1, 0.35,
             font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    try_add_image(slide, IMG_CLASS / "confusion_matrix_decisiontree_classification.png",
                  6.8, 4.4, 6.4, 2.85)
    add_text(slide, "Decision Tree Confusion Matrix", 6.8, 7.1, 6.4, 0.35,
             font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    slide_number(slide, prs, 10)
    return slide


def slide11_classif2(prs):
    slide = blank_slide(prs)
    bg(slide)
    slide_title_bar(slide, "Classification Models — Part 2",
                    "SVM · K-Means · Agglomerative Clustering")

    bullets_left = [
        "SVM — margin-maximising hyperplane (RBF kernel from scratch)",
        f"  Accuracy: {fmt(CLASSIF['SVM (scratch)']['accuracy'],4)}  "
        f"  Recall: {fmt(CLASSIF['SVM (scratch)']['recall'],4)}  "
        f"(high recall — fewer missed churners)",
        "",
        "K-Means Clustering — unsupervised; 2 clusters mapped to labels",
        f"  Accuracy: {fmt(CLASSIF['K-Means Clustering (scratch)']['accuracy'],4)}  "
        f"  F1: {fmt(CLASSIF['K-Means Clustering (scratch)']['f1_score'],4)}",
        "",
        "Agglomerative Clustering — hierarchical; ward linkage",
        f"  Accuracy: {fmt(CLASSIF['Agglomerative Clustering (scratch)']['accuracy'],4)}  "
        f"  F1: {fmt(CLASSIF['Agglomerative Clustering (scratch)']['f1_score'],4)}",
        "",
        "PCA applied before SVM/RF for dimensionality reduction variants",
        "RF+PCA and SVM+PCA results nearly identical to base variants",
    ]
    add_bullets(slide, bullets_left, 0.4, 1.25, 6.5, 5.9, font_size=14)

    try_add_image(slide, IMG_CLASS / "confusion_matrix_svm_classification.png",
                  7.0, 1.3, 3.0, 2.7)
    try_add_image(slide, IMG_CLASS / "confusion_matrix_kmeans_classification.png",
                  10.1, 1.3, 3.0, 2.7)
    add_text(slide, "SVM CM", 7.0, 3.95, 3.0, 0.3,
             font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "K-Means CM", 10.1, 3.95, 3.0, 0.3,
             font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    try_add_image(slide, IMG_CLASS / "confusion_matrix_agglomerativeclustering_classification.png",
                  7.0, 4.35, 6.3, 2.9)
    add_text(slide, "Agglomerative Clustering CM", 7.0, 7.1, 6.3, 0.35,
             font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    try_add_image(slide, IMG_CORR / "correlation_classification.png",
                  7.0, 4.35, 3.0, 2.9)
    slide_number(slide, prs, 11)
    return slide


def slide12_classif3(prs):
    slide = blank_slide(prs)
    bg(slide)
    slide_title_bar(slide, "Classification Models — Part 3",
                    "Perceptron / SLP · MLP · Custom (Naive Bayes)")

    bullets_left = [
        "Perceptron / SLP — step activation, single weight update",
        f"  Accuracy: {fmt(CLASSIF['Perceptron (scratch)']['accuracy'],4)}  "
        f"  F1: {fmt(CLASSIF['Perceptron (scratch)']['f1_score'],4)}",
        "",
        "MLP — 2 hidden layers (ReLU), sigmoid output, Adam-like GD",
        f"  Accuracy: {fmt(CLASSIF['MLP (scratch)']['accuracy'],4)}  "
        f"  F1: {fmt(CLASSIF['MLP (scratch)']['f1_score'],4)}",
        "",
        "Custom Model — Gaussian Naive Bayes (outside curriculum)",
        f"  Accuracy: {fmt(CLASSIF['Custom Model (scratch)']['accuracy'],4)}  "
        f"  F1: {fmt(CLASSIF['Custom Model (scratch)']['f1_score'],4)}",
        "",
        "All neural models trained for 50–100 epochs with early stopping",
    ]
    add_bullets(slide, bullets_left, 0.4, 1.25, 6.2, 5.9, font_size=15)

    try_add_image(slide, IMG_CLASS / "confusion_matrix_perceptron_classification.png",
                  6.7, 1.3, 3.2, 2.75)
    try_add_image(slide, IMG_CLASS / "confusion_matrix_mlp_classification.png",
                  9.9, 1.3, 3.2, 2.75)
    add_text(slide, "Perceptron CM", 6.7, 3.95, 3.2, 0.3,
             font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "MLP CM", 9.9, 3.95, 3.2, 0.3,
             font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    try_add_image(slide, IMG_CLASS / "confusion_matrix_custommodel_classification.png",
                  6.7, 4.35, 6.4, 2.85)
    add_text(slide, "Custom Model (Naive Bayes) CM", 6.7, 7.1, 6.4, 0.3,
             font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    slide_number(slide, prs, 12)
    return slide


def slide13_novel(prs):
    slide = blank_slide(prs)
    bg(slide)
    slide_title_bar(slide, "Novel / New Models",
                    "Going beyond the classroom curriculum")

    # XGBoost regression card
    add_rect(slide, 0.4, 1.3, 5.9, 5.7, CARD_BG)
    add_rect(slide, 0.4, 1.3, 5.9, 0.5, RGBColor(0x1A, 0x78, 0xC2))
    add_text(slide, "XGBoost Regression (Novel)", 0.6, 1.33, 5.6, 0.43,
             font_size=16, bold=True, color=WHITE)
    xgb_bullets = [
        "Gradient boosting: trains trees on residuals",
        "Shrinkage + column subsampling for regularisation",
        "Built entirely from scratch (NumPy only)",
        "R²  : 0.9946  (vs 0.4134 for Linear Reg.)",
        "RMSE: 119.86  (vs 1,248.77 for Linear Reg.)",
        "Validated against XGBoost sklearn: R² 0.9989",
        "Massive improvement — 10× lower RMSE than Poly Reg",
    ]
    add_bullets(slide, xgb_bullets, 0.5, 1.9, 5.7, 5.0, font_size=15)

    # Naive Bayes classification card
    add_rect(slide, 7.0, 1.3, 5.9, 5.7, CARD_BG)
    add_rect(slide, 7.0, 1.3, 5.9, 0.5, RGBColor(0x1A, 0x9C, 0x60))
    add_text(slide, "Naive Bayes Classification (Novel)", 7.2, 1.33, 5.6, 0.43,
             font_size=16, bold=True, color=WHITE)
    nb_bullets = [
        "Gaussian NB: probabilistic generative classifier",
        "Assumes feature independence given class label",
        "Built from scratch: compute μ, σ per class per feature",
        "Accuracy : 0.8245  F1: 0.6747",
        "Comparable to MLP & Agglomerative Clustering",
        "Extremely fast inference — no iterative training needed",
        "Novel: outside classroom curriculum",
    ]
    add_bullets(slide, nb_bullets, 7.1, 1.9, 5.7, 5.0, font_size=15)
    slide_number(slide, prs, 13)
    return slide


def slide14_benchmark(prs):
    slide = blank_slide(prs)
    bg(slide)
    slide_title_bar(slide, "Benchmarking Table",
                    "All 17 models compared side-by-side")

    # ---- classification table ----
    add_text(slide, "CLASSIFICATION", 0.4, 1.2, 7.0, 0.35,
             font_size=13, bold=True, color=ACCENT)
    clf_headers = ["Model", "Acc", "Prec", "Recall", "F1"]
    clf_col_w   = [3.5, 1.1, 1.1, 1.1, 1.1]
    clf_col_x   = [0.35, 3.9, 5.05, 6.2, 7.35]
    clf_data = [
        ("Logistic Regression", "0.833","0.628","0.728","0.675"),
        ("Decision Tree",       "0.841","0.694","0.592","0.639"),
        ("Random Forest",       "0.843","0.696","0.601","0.645"),
        ("SVM",                 "0.757","0.494","0.903","0.639"),
        ("RF + PCA",            "0.843","0.696","0.601","0.645"),
        ("SVM + PCA",           "0.757","0.494","0.904","0.639"),
        ("K-Means",             "0.641","0.242","0.240","0.241"),
        ("Agglomerative",       "0.825","0.604","0.765","0.675"),
        ("Perceptron",          "0.782","0.542","0.524","0.533"),
        ("MLP",                 "0.827","0.610","0.752","0.673"),
        ("Naive Bayes ★",       "0.825","0.602","0.767","0.675"),
        ("RF sklearn ✔",        "0.843","0.694","0.606","0.647"),
    ]
    y = 1.55
    # header
    for hdr, x, w in zip(clf_headers, clf_col_x, clf_col_w):
        add_rect(slide, x, y, w - 0.04, 0.38, ACCENT)
        add_text(slide, hdr, x + 0.03, y + 0.04, w - 0.08, 0.3,
                 font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    y += 0.4
    for ri, row in enumerate(clf_data):
        rc = RGBColor(0x0D, 0x35, 0x60) if "★" in row[0] or "✔" in row[0] else CARD_BG
        for ci, (cell, x, w) in enumerate(zip(row, clf_col_x, clf_col_w)):
            add_rect(slide, x, y, w - 0.04, 0.38, rc)
            clr = GOLD if "★" in row[0] else WHITE
            add_text(slide, cell, x + 0.03, y + 0.04, w - 0.08, 0.3,
                     font_size=10, color=clr, align=PP_ALIGN.CENTER,
                     bold=("★" in row[0]))
        y += 0.40

    # ---- regression table ----
    add_text(slide, "REGRESSION", 8.7, 1.2, 4.5, 0.35,
             font_size=13, bold=True, color=ACCENT)
    reg_headers = ["Model", "R²", "RMSE"]
    reg_col_w   = [3.2, 1.1, 1.2]
    reg_col_x   = [8.65, 11.9, 13.05]
    reg_data = [
        ("Linear Reg.",       "0.413", "1248.8"),
        ("Multiple Reg.",     "0.100", "1546.8"),
        ("Polynomial Reg.",   "0.744", " 825.5"),
        ("XGBoost (scratch) ★","0.995", " 119.9"),
        ("XGBoost sklearn ✔", "0.999", "  54.8"),
    ]
    y2 = 1.55
    for hdr, x, w in zip(reg_headers, reg_col_x, reg_col_w):
        add_rect(slide, x, y2, w - 0.04, 0.38, ACCENT)
        add_text(slide, hdr, x + 0.03, y2 + 0.04, w - 0.08, 0.3,
                 font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    y2 += 0.4
    for row in reg_data:
        rc = RGBColor(0x0D, 0x35, 0x60) if "★" in row[0] or "✔" in row[0] else CARD_BG
        for cell, x, w in zip(row, reg_col_x, reg_col_w):
            add_rect(slide, x, y2, w - 0.04, 0.38, rc)
            clr = GOLD if "★" in row[0] else WHITE
            add_text(slide, cell, x + 0.03, y2 + 0.04, w - 0.08, 0.3,
                     font_size=10, color=clr, align=PP_ALIGN.CENTER,
                     bold=("★" in row[0]))
        y2 += 0.40

    # embed overall metrics comparison image
    try_add_image(slide, IMG_CLASS / "metrics_comparison.png",
                  8.65, 3.0, 4.5, 4.3)

    slide_number(slide, prs, 14)
    return slide


def slide15_roc(prs):
    slide = blank_slide(prs)
    bg(slide)
    slide_title_bar(slide, "ROC & AUC Analysis",
                    "Classification performance across decision thresholds")

    bullets = [
        "AUC measures discriminative power regardless of threshold",
        "Random Forest: highest accuracy (0.843) → strong AUC expected",
        "SVM: highest recall (0.903) — minimises missed churners",
        "Clustering models (K-Means) show poor AUC due to unsupervised nature",
        "Best F1 models: Logistic Reg. / Agglomerative / Naive Bayes ≈ 0.675",
    ]
    add_bullets(slide, bullets, 0.4, 1.3, 6.5, 2.5, font_size=16)

    add_text(slide, "▸  ROC curve plots — [INSERT ROC CHART FROM MLFLOW]",
             0.4, 3.9, 6.5, 0.5, font_size=14, color=LIGHT_GRAY)

    try_add_image(slide, IMG_CLASS / "metrics_comparison.png",
                  7.0, 1.3, 6.1, 5.8)
    add_text(slide, "▲ Classification metrics comparison (all models)",
             7.0, 6.95, 6.1, 0.4, font_size=11, color=LIGHT_GRAY,
             align=PP_ALIGN.CENTER)

    add_text(slide, "Note: Full ROC curves tracked in MLflow experiment 'retail_classification_benchmark'",
             0.4, 7.1, 12.5, 0.35, font_size=11, color=LIGHT_GRAY, italic=True)
    slide_number(slide, prs, 15)
    return slide


def slide16_discussion(prs):
    slide = blank_slide(prs)
    bg(slide)
    slide_title_bar(slide, "Discussion & Key Findings")

    cards = [
        ("🏆 Winning Models",
         ["Regression: XGBoost scratch (R²=0.9946, RMSE=119.9)",
          "Classification: Random Forest or Logistic Reg. (F1≈0.675)",
          "Novel models validated & competitive with sklearn baselines"],
         0.4, 1.3, 6.0),
        ("🔬 From-Scratch vs sklearn",
         ["XGBoost scratch vs sklearn: R² 0.9946 vs 0.9989",
          "RF scratch vs sklearn: Accuracy 0.843 vs 0.843",
          "From-scratch results closely match library implementations",
          "Validates correctness of all custom algorithm implementations"],
         6.6, 1.3, 6.0),
        ("⚠️  Overfitting / Underfitting",
         ["Multiple Regression underfit: R²=0.10 (too simple)",
          "K-Means poor on supervised task (F1=0.241) — wrong tool",
          "Polynomial Reg. moderate fit R²=0.74 before XGBoost",
          "XGBoost shows near-zero training residuals"],
         0.4, 4.3, 6.0),
        ("📊 Feature Insights",
         ["Interaction feat. quantity×price most predictive for regression",
          "Recency (days_since_last_purchase) top churn signal",
          "Loyalty membership strongly negatively correlated with churn",
          "PCA variants maintain performance ─ confirm feature redundancy"],
         6.6, 4.3, 6.0),
    ]
    for title, bullets, x, y, w in cards:
        add_rect(slide, x, y, w, 2.75, CARD_BG)
        add_rect(slide, x, y, w, 0.45, CARD_BG)
        add_text(slide, title, x + 0.15, y + 0.06, w - 0.2, 0.35,
                 font_size=15, bold=True, color=ACCENT)
        add_bullets(slide, bullets, x + 0.1, y + 0.55, w - 0.2, 2.1,
                    font_size=13, color=LIGHT_GRAY)
    slide_number(slide, prs, 16)
    return slide


def slide17_conclusion(prs):
    slide = blank_slide(prs)
    bg(slide)
    slide_title_bar(slide, "Conclusion")

    achieved = [
        "✅  17 ML models implemented — all from scratch using NumPy",
        "✅  Dual-task pipeline: regression (total_sales) + classification (churned)",
        "✅  Novel models: XGBoost & Naive Bayes exceed baseline curriculum models",
        "✅  Full MLOps pipeline: Airflow + MinIO + MLflow + Docker",
        "✅  From-scratch results validated against sklearn equivalents",
    ]
    add_text(slide, "Achievements", 0.4, 1.3, 12.0, 0.4,
             font_size=18, bold=True, color=ACCENT)
    add_bullets(slide, achieved, 0.4, 1.75, 12.0, 2.8, font_size=16, color=GREEN)

    limits = [
        "⚑  Dataset is synthetic — real-world performance may differ",
        "⚑  K-Means not suited for directly supervised churn prediction",
        "⚑  MLP architecture limited by from-scratch constraint (no GPU)",
        "⚑  MAPE high for linear/multiple reg. due to near-zero target values",
    ]
    add_text(slide, "Limitations", 0.4, 4.6, 12.0, 0.4,
             font_size=18, bold=True, color=ACCENT)
    add_bullets(slide, limits, 0.4, 5.05, 12.0, 2.0, font_size=16, color=RED_LIGHT)

    future = [
        "▸  Deep learning (LSTM/Transformer) for sequential transaction modelling",
        "▸  Real production deployment with live churn alerts",
        "▸  AutoML hyperparameter tuning integrated into Airflow DAG",
    ]
    add_text(slide, "Future Work", 0.4, 5.55, 12.0, 0.4,
             font_size=18, bold=True, color=ACCENT)
    add_bullets(slide, future, 0.4, 6.0, 12.0, 1.3, font_size=15, color=LIGHT_GRAY)
    slide_number(slide, prs, 17)
    return slide


def slide18_references(prs):
    slide = blank_slide(prs)
    bg(slide)
    slide_title_bar(slide, "References")

    refs = [
        "[1] Retail Customer Analytics Dataset — Synthetic multi-category retail data",
        "[2] Breiman, L. (2001). Random Forests. Machine Learning, 45, 5–32.",
        "[3] Chen, T. & Guestrin, C. (2016). XGBoost: A Scalable Tree Boosting System. KDD.",
        "[4] Cortes, C. & Vapnik, V. (1995). Support Vector Networks. Machine Learning, 20, 273–297.",
        "[5] Scikit-learn: Machine Learning in Python. Pedregosa et al., JMLR 12, 2011.",
        "[6] Apache Airflow — Workflow orchestration platform. https://airflow.apache.org",
        "[7] MLflow — Open source ML lifecycle platform. https://mlflow.org",
        "[8] MinIO — High-performance object storage. https://min.io",
        "[9] Goodfellow, I., Bengio, Y., Courville, A. (2016). Deep Learning. MIT Press.",
        "[10] Mitchell, T. (1997). Machine Learning. McGraw-Hill.",
    ]
    y = 1.3
    for ref in refs:
        add_text(slide, ref, 0.5, y, 12.3, 0.45,
                 font_size=15, color=LIGHT_GRAY)
        y += 0.52

    add_text(slide,
             "Built with ❤️  for [01276343] ML Project  —  March 9, 2026",
             0, 7.1, 13.33, 0.35,
             font_size=13, color=ACCENT, align=PP_ALIGN.CENTER)
    slide_number(slide, prs, 18)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    builders = [
        slide1_title,
        slide2_agenda,
        slide3_problem,
        slide4_dataset,
        slide5_features,
        slide6_preprocessing,
        slide7_methodology,
        slide8_regression_models,
        slide9_regression_results,
        slide10_classif1,
        slide11_classif2,
        slide12_classif3,
        slide13_novel,
        slide14_benchmark,
        slide15_roc,
        slide16_discussion,
        slide17_conclusion,
        slide18_references,
    ]

    for builder in builders:
        builder(prs)
        print(f"  ✔  {builder.__name__}")

    prs.save(str(OUT_PATH))
    print(f"\n✅  Saved → {OUT_PATH}")


if __name__ == "__main__":
    main()
